from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# 创建新的PPT
prs = Presentation()

# 设置幻灯片宽高
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 添加标题幻灯片
slide_layout = prs.slide_layouts[5]  # 空白布局
slide = prs.slides.add_slide(slide_layout)

# 添加标题
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(1))
title_frame = title_shape.text_frame
title_frame.text = "FDA 对光学轮廓感应设备的指导意见摘要"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x49, 0x7C)

# 内容框架
content = [
    ("🎯 设备概述", [
        "用途：测量脊柱、足部等解剖标志",
        "应用：姿势监测、脊柱异常检测（侧弯/驼背/前凸）",
        "技术：摄像头、光学扫描 + 软件分析"
    ]),
    ("⚠️ 风险识别（FDA）", [
        "设备故障或误操作 → 误诊/病情恶化",
        "需确认是否补充其他潜在风险"
    ]),
    ("🧪 FDA 分类建议", [
        "分类提议：Class I（仅需通用控制）",
        "理由：安全有效，无需特别控制",
        "征求反馈：是否认同该分类？"
    ]),
    ("❓ 关键问题", [
        "是否认同风险识别与分类？",
        "是否有遗漏风险？是否需补充？"
    ])
]

# 添加内容文本框
left = Inches(0.5)
top = Inches(1.3)
width = Inches(12.3)
height = Inches(5.8)
content_box = slide.shapes.add_textbox(left, top, width, height)
content_frame = content_box.text_frame
content_frame.word_wrap = True

# 填充内容
for section_title, bullet_points in content:
    p = content_frame.add_paragraph()
    p.text = section_title
    p.font.bold = True
    p.font.size = Pt(20)
    p.space_before = Pt(10)
    for item in bullet_points:
        bullet = content_frame.add_paragraph()
        bullet.text = item
        bullet.level = 1
        bullet.font.size = Pt(16)

# 保存PPT
ppt_path = "/mnt/data/FDA_Optical_Device_Summary.pptx"
prs.save(ppt_path)

ppt_path
